import streamlit as st
from dotenv import load_dotenv
from PyPDF2 import PdfReader
import nltk
import re
import slate3k as slate
from pdf2image import convert_from_path
from nltk.corpus import stopwords
from nltk.tokenize import word_tokenize, sent_tokenize
from nltk.stem.snowball import SnowballStemmer
from PIL import Image

from langchain.text_splitter import CharacterTextSplitter
from langchain.embeddings.openai import OpenAIEmbeddings
from langchain.vectorstores import FAISS
from langchain.chains.question_answering import load_qa_chain
from langchain.llms import OpenAI
from langchain.callbacks import get_openai_callback

from pptx import Presentation
from pptx.util import Inches
import base64
import random
import openai
import io

nltk.download("stopwords")
nltk.download("punkt")


openai.api_key = "your api key here"

Prompt = """
Write a presentation/powerpoint about the user's topic. You only answer with the presentation. Follow the structure of the example.
Notice
- You do all the presentation text for the user.
- You write the texts no longer than 250 characters!
- You make very short titles!
- You make the presentation easy to understand.
- The presentation has a table of contents.
- The presentation has a summary.
- At least 4 slides.
- At most 10 slides

Example! - Stick to this formatting exactly!
#Title: TITLE OF THE PRESENTATION

#Slide: 1
#Header: table of contents
#Content: 1. CONTENT OF THIS POWERPOINT
2. CONTENTS OF THIS POWERPOINT
3. CONTENT OF THIS POWERPOINT
...

#Slide: 2
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 3
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 4
#Header: TITLE OF SLIDE
#Content: CONTENT OF THE SLIDE

#Slide: 5
#Headers: summary
#Content: CONTENT OF THE SUMMARY

#Slide: END
"""


def create_ppt_text(input):
    response = openai.ChatCompletion.create(
        model="gpt-3.5-turbo",
        messages=[
            {"role": "system", "content": Prompt},
            {"role": "user", "content": f"The user wants a presentation about {input}"},
        ],
        temperature=0.5,
    )

    return response['choices'][0]['message']['content']


def create_ppt(input_text, ppt_name):
    prs = Presentation()
    slide_count = 0
    header = ""
    content = ""
    last_slide_layout_index = -1
    first_time = True
    lines = input_text.split("\n")
    for line in lines:
        if line.startswith('#Title:'):
            header = line.replace('#Title:', '').strip()
            slide = prs.slides.add_slide(prs.slide_layouts[0])
            title = slide.shapes.title
            title.text = header
            body_shape = slide.shapes.placeholders[1]
        elif line.startswith('#Slide:'):
            if slide_count > 0:
                slide = prs.slides.add_slide(prs.slide_layouts[slide_layout_index])
                title = slide.shapes.title
                title.text = header
                body_shape = slide.shapes.placeholders[slide_placeholder_index]
                tf = body_shape.text_frame
                tf.text = content
            content = ""
            slide_count += 1
            slide_layout_index = last_slide_layout_index
            layout_indices = [1, 7, 8]
            while slide_layout_index == last_slide_layout_index:
                if first_time:
                    slide_layout_index = 1
                    slide_placeholder_index = 1
                    first_time = False
                    break
                slide_layout_index = random.choice(layout_indices)  # Select random slide index
                if slide_layout_index == 8:
                    slide_placeholder_index = 2
                else:
                    slide_placeholder_index = 1
            last_slide_layout_index = slide_layout_index
        elif line.startswith('#Header:'):
            header = line.replace('#Header:', '').strip()
        elif line.startswith('#Content:'):
            content = line.replace('#Content:', '').strip()
            next_line_index = lines.index(line) + 1
            while next_line_index < len(lines) and not lines[next_line_index].startswith('#'):
                content += '\n' + lines[next_line_index].strip()
                next_line_index += 1

    ppt_bytes = ppt_to_bytes(prs)
    return ppt_bytes


def ppt_to_bytes(presentation):
    ppt_bytes = io.BytesIO()
    presentation.save(ppt_bytes)
    ppt_bytes.seek(0)
    return base64.b64encode(ppt_bytes.read()).decode("utf-8")



def main():
    load_dotenv()
    st.set_page_config(page_title="StudyBuddy", layout="wide")
    st.title("StudyBuddy")

    css = """
        <style>
        .title-container {
            background-color: white;
            padding: 10px;
            border: 1px solid white;
            border-radius: 5px;
            margin-bottom: 20px;
        }
        </style>
        """

    # Render the CSS styling
    st.markdown(css, unsafe_allow_html=True)
    st.markdown('<div class="title-container">', unsafe_allow_html=True)
    st.markdown("</div>", unsafe_allow_html=True)


    # Tab selection
    tab_options = ["Ask PDF", "PDF Summarizer", "AI Presentation Maker"]  # Add the third tab option
    selected_tab = st.selectbox("Select a tab", tab_options)


    if selected_tab == "Ask PDF":
        ask_pdf_tab()
    elif selected_tab == "PDF Summarizer":
        pdf_summarizer_tab()

    elif selected_tab == "AI Presentation Maker":  # Add the condition for the new tab
        ai_presentation_maker_tab()  # Call the function for the new tab


def ask_pdf_tab():
    st.title("Ask your PDF 💬")

    # Upload file
    pdf = st.file_uploader("Upload your PDF", type="pdf")

    # Extract the text
    if pdf is not None:
        pdf_reader = PdfReader(pdf)
        text = ""
        for page in pdf_reader.pages:
            text += page.extract_text()

        # Split into chunks
        text_splitter = CharacterTextSplitter(
            separator="\n",
            chunk_size=1000,
            chunk_overlap=200,
            length_function=len
        )
        chunks = text_splitter.split_text(text)

        # Create embeddings
        embeddings = OpenAIEmbeddings()
        knowledge_base = FAISS.from_texts(chunks, embeddings)

        # Show user input
        user_question = st.text_input("Ask a question about your PDF:")
        if user_question:
            docs = knowledge_base.similarity_search(user_question)

            llm = OpenAI()
            chain = load_qa_chain(llm, chain_type="stuff")
            with get_openai_callback() as cb:
                response = chain.run(input_documents=docs, question=user_question)
                print(cb)

            st.write(response)


def pdf_summarizer_tab():
    st.title("PDF Summarizer")

    # Upload PDF file
    uploaded_file = st.file_uploader("Upload a PDF file", type="pdf")
    if uploaded_file is not None:
        file_name = uploaded_file.name
        with open(file_name, "wb") as file:
            file.write(uploaded_file.getbuffer())

        # Text extraction option
        extraction_option = st.radio("Choose extraction option", ("Text",))
        if extraction_option == "Text":
            text = extract_text(file_name)

        # Summarize text
        summary = summarize(text)

        # Display summary
        st.header("Summary")
        st.markdown(summary)


def extract_text(file):
    pdfFileObj = open(file, "rb")
    pdfPages = slate.PDF(pdfFileObj)

    # Extract text from PDF file
    text = ""
    for page in pdfPages:
        text += page
    return text


def summarize(text):
    # Process text by removing numbers and unrecognized punctuation
    processedText = re.sub("’", "'", text)
    processedText = re.sub("[^a-zA-Z' ]+", " ", processedText)
    stopWords = set(stopwords.words("english"))
    words = word_tokenize(processedText)

    # Normalize words with Porter stemming and build word frequency table
    stemmer = SnowballStemmer("english", ignore_stopwords=True)
    freqTable = dict()
    for word in words:
        word = word.lower()
        if word in stopWords:
            continue
        elif stemmer.stem(word) in freqTable:
            freqTable[stemmer.stem(word)] += 1
        else:
            freqTable[stemmer.stem(word)] = 1

    # Normalize every sentence in the text
    sentences = sent_tokenize(text)
    stemmedSentences = []
    sentenceValue = dict()
    for sentence in sentences:
        stemmedSentence = []
        for word in sentence.lower().split():
            stemmedSentence.append(stemmer.stem(word))
        stemmedSentences.append(stemmedSentence)

    # Calculate value of every normalized sentence based on word frequency table
    # [:12] helps to save space
    for num in range(len(stemmedSentences)):
        for wordValue in freqTable:
            if wordValue in stemmedSentences[num]:
                if sentences[num][:12] in sentenceValue:
                    sentenceValue[sentences[num][:12]] += freqTable.get(wordValue)
                else:
                    sentenceValue[sentences[num][:12]] = freqTable.get(wordValue)

    # Determine average value of a sentence in the text
    sumValues = 0
    for sentence in sentenceValue:
        sumValues += sentenceValue.get(sentence)

    average = int(sumValues / len(sentenceValue))

    # Create summary of text using sentences that exceed the average value by some factor
    # This factor can be adjusted to reduce/expand the length of the summary
    summary = ""
    for sentence in sentences:
        if sentence[:12] in sentenceValue and sentenceValue[sentence[:12]] > (3.0 * average):
            summary += " " + " ".join(sentence.split())

    # Process the text in summary and return it
    summary = re.sub("’", "'", summary)
    summary = re.sub("[^a-zA-Z0-9'\"():;,.!?— ]+", " ", summary)
    return summary


def ai_presentation_maker_tab():  # Add the function for the new tab
    st.title("AI Presentation Maker")

    user_text = st.text_input("Enter the topic for the presentation")
    if st.button("Create Presentation"):
        input_string = user_text.strip()
        input_string = re.sub(r'[^\w\s.\-\(\)]', '', input_string)
        input_string = input_string.replace("\n", "")
        ppt_bytes = create_ppt(create_ppt_text(input_string), input_string)
        st.success("Presentation created successfully!")
        st.markdown(
            f'<a href="data:application/vnd.openxmlformats-officedocument.presentationml.presentation;base64,{ppt_bytes}" download="{input_string}.pptx">Download Presentation</a>',
            unsafe_allow_html=True)

if __name__ == '__main__':
    main()